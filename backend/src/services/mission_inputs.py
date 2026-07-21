"""Seal uploaded text as immutable Mission inputs and project it into chat context."""

from __future__ import annotations

import hashlib
import mimetypes
import os
import secrets
import stat
import zipfile
from collections.abc import Iterable
from dataclasses import dataclass
from datetime import date, datetime, time
from io import BytesIO
from pathlib import Path, PurePosixPath
from typing import Any

from src.application.results import ThreadTurnAttachment
from src.config import get_settings
from src.contracts.archive_filename import recover_legacy_zip_filename
from src.contracts.mission_input import (
    MAX_MISSION_INPUTS,
    MissionInputContext,
    MissionInputManifest,
)
from src.services.path_safety import normalize_path_component
from src.services.workspace_uploads import is_pdf_upload

_THREAD_VIRTUAL_ROOT = "/mnt/user-data"
_MAX_SOURCE_BYTES = 100 * 1024 * 1024
_MAX_TEXT_BYTES = 8 * 1024 * 1024
_MAX_PROMPT_INPUTS = 8
_MAX_EXCERPT_CHARS = 12_000
_MAX_OOXML_FILES = 4096
_MAX_OOXML_UNCOMPRESSED_BYTES = 128 * 1024 * 1024
_MAX_OOXML_PART_BYTES = 32 * 1024 * 1024
_MAX_SPREADSHEET_ROWS = 200_000
_MAX_SPREADSHEET_COLUMNS = 4096
_MAX_PRESENTATION_SLIDES = 2000
_MAX_ARCHIVE_FILES = 128
_MAX_ARCHIVE_INPUTS = MAX_MISSION_INPUTS
_MAX_ARCHIVE_UNCOMPRESSED_BYTES = 256 * 1024 * 1024
_MAX_ARCHIVE_MEMBER_BYTES = 100 * 1024 * 1024
_MAX_ARCHIVE_COMPRESSION_RATIO = 200
_TEXT_SUFFIXES = frozenset(
    {
        ".bib",
        ".csv",
        ".json",
        ".md",
        ".markdown",
        ".py",
        ".r",
        ".tex",
        ".tsv",
        ".txt",
        ".xml",
        ".yaml",
        ".yml",
    }
)


@dataclass(frozen=True, slots=True)
class MissionInputPreparation:
    manifests: tuple[MissionInputManifest, ...]
    contexts: tuple[MissionInputContext, ...]


class MissionInputStore:
    """Private workspace-scoped object store for immutable extracted text."""

    def __init__(self, root: Path | None = None) -> None:
        self.root = Path(root or get_settings().thread_data_root)

    def put_text(
        self,
        *,
        workspace_id: str,
        thread_id: str,
        filename: str,
        mime_type: str | None,
        extractor: str,
        text: str,
        source_content_hash: str,
        source_size_bytes: int,
        container_filename: str | None = None,
        container_content_hash: str | None = None,
        member_path: str | None = None,
    ) -> MissionInputManifest:
        normalized = _normalize_text(text)
        payload = normalized.encode("utf-8")
        if not payload or len(payload) > _MAX_TEXT_BYTES:
            raise ValueError("mission input extracted text is empty or exceeds the hard limit")
        digest = hashlib.sha256(payload).hexdigest()
        manifest = MissionInputManifest(
            input_ref=f"mission-input:{digest}",
            workspace_id=workspace_id,
            thread_id=thread_id,
            filename=Path(filename).name,
            mime_type=mime_type,
            extractor=extractor,
            content_hash=f"sha256:{digest}",
            source_content_hash=source_content_hash,
            source_size_bytes=source_size_bytes,
            text_size_bytes=len(payload),
            text_chars=len(normalized),
            container_filename=container_filename,
            container_content_hash=container_content_hash,
            member_path=member_path,
        )
        target = self._object_path(manifest)
        _ensure_private_directory(target.parent)
        if target.exists():
            self._verify_file(target, manifest)
        else:
            _atomic_write_once(target, payload)
            self._verify_file(target, manifest)
        return manifest

    def read_text(
        self,
        manifest: MissionInputManifest,
        *,
        workspace_id: str,
        thread_id: str | None = None,
        offset: int = 0,
        max_chars: int = 65_536,
    ) -> tuple[str, bool]:
        if manifest.workspace_id != workspace_id:
            raise PermissionError("mission input workspace mismatch")
        if thread_id is not None and manifest.thread_id != thread_id:
            raise PermissionError("mission input thread mismatch")
        target = self._object_path(manifest)
        payload = self._verify_file(target, manifest)
        text = payload.decode("utf-8")
        if offset > len(text):
            return "", False
        end = min(len(text), offset + max_chars)
        return text[offset:end], end < len(text)

    def _object_path(self, manifest: MissionInputManifest) -> Path:
        digest = manifest.input_ref.removeprefix("mission-input:")
        thread = normalize_path_component(manifest.thread_id)
        return self.root / thread / "mission-inputs" / digest[:2] / digest / "content.txt"

    @staticmethod
    def _verify_file(path: Path, manifest: MissionInputManifest) -> bytes:
        if path.is_symlink():
            raise PermissionError("mission input object cannot be a symlink")
        try:
            payload = path.read_bytes()
        except FileNotFoundError as exc:
            raise LookupError("mission input object was not found") from exc
        digest = hashlib.sha256(payload).hexdigest()
        if f"sha256:{digest}" != manifest.content_hash or len(payload) != manifest.text_size_bytes or len(payload.decode("utf-8")) != manifest.text_chars:
            raise ValueError("mission input object integrity check failed")
        return payload


class MissionInputService:
    """Convert thread uploads into durable, typed inputs without trusting client paths."""

    def __init__(
        self,
        *,
        store: MissionInputStore | None = None,
        thread_data_root: Path | None = None,
    ) -> None:
        settings = get_settings()
        self.store = store or MissionInputStore(settings.thread_data_root)
        self.thread_data_root = Path(thread_data_root or settings.thread_data_root)

    def prepare(
        self,
        *,
        workspace_id: str,
        thread_id: str,
        attachments: tuple[ThreadTurnAttachment, ...],
    ) -> MissionInputPreparation:
        manifests: list[MissionInputManifest] = []
        contexts: list[MissionInputContext] = []
        for attachment in attachments:
            prepared = self._prepare_attachment(
                workspace_id=workspace_id,
                thread_id=thread_id,
                attachment=attachment,
            )
            for context in prepared.contexts:
                if len(contexts) < MAX_MISSION_INPUTS:
                    contexts.append(context)
            for manifest in prepared.manifests:
                if all(item.input_ref != manifest.input_ref for item in manifests):
                    manifests.append(manifest)
        if len(manifests) > MAX_MISSION_INPUTS:
            raise ValueError(f"one turn may prepare at most {MAX_MISSION_INPUTS} Mission inputs")
        return MissionInputPreparation(manifests=tuple(manifests), contexts=tuple(contexts))

    def _prepare_attachment(
        self,
        *,
        workspace_id: str,
        thread_id: str,
        attachment: ThreadTurnAttachment,
    ) -> MissionInputPreparation:
        if not _is_zip_upload(attachment.name, attachment.content_type):
            manifest, context = self._prepare_one(
                workspace_id=workspace_id,
                thread_id=thread_id,
                attachment=attachment,
            )
            return MissionInputPreparation(
                manifests=(() if manifest is None else (manifest,)),
                contexts=(context,),
            )
        base = {
            "name": attachment.name,
            "content_type": attachment.content_type,
            "size_bytes": attachment.size_bytes,
        }
        try:
            source_path = self._resolve_thread_path(thread_id, attachment.path)
            source = source_path.read_bytes()
        except (OSError, ValueError, PermissionError):
            return MissionInputPreparation(
                manifests=(),
                contexts=(MissionInputContext(**base, status="unreadable", detail="压缩包不可读取，请重新上传。"),),
            )
        if not source or len(source) > _MAX_SOURCE_BYTES:
            return MissionInputPreparation(
                manifests=(),
                contexts=(MissionInputContext(**base, status="unreadable", detail="压缩包为空或超过可处理大小。"),),
            )
        try:
            members = _read_safe_zip_members(source)
        except ValueError as exc:
            return MissionInputPreparation(
                manifests=(),
                contexts=(MissionInputContext(**base, status="unreadable", detail=str(exc)),),
            )
        container_hash = f"sha256:{hashlib.sha256(source).hexdigest()}"
        manifests: list[MissionInputManifest] = []
        contexts: list[MissionInputContext] = []
        for member_path, member_content in members:
            member_mime = mimetypes.guess_type(member_path)[0]
            manifest, context = self._prepare_source(
                workspace_id=workspace_id,
                thread_id=thread_id,
                name=f"{attachment.name}/{member_path}",
                content_type=member_mime,
                source=member_content,
                container_filename=attachment.name,
                container_content_hash=container_hash,
                member_path=member_path,
            )
            if manifest is None:
                continue
            if all(item.input_ref != manifest.input_ref for item in manifests):
                manifests.append(manifest)
                contexts.append(context)
            if len(manifests) >= _MAX_ARCHIVE_INPUTS:
                break
        if not manifests:
            return MissionInputPreparation(
                manifests=(),
                contexts=(MissionInputContext(**base, status="unreadable", detail="压缩包中没有可读取的 PDF、Office 或文本文件。"),),
            )
        total_readable = sum(
            1
            for member_path, _ in members
            if _is_supported_member(member_path, mimetypes.guess_type(member_path)[0])
        )
        if total_readable > len(manifests):
            first = contexts[0]
            contexts[0] = first.model_copy(
                update={
                    "detail": f"已读取压缩包成员；受任务输入上限约束，本次读取 {len(manifests)}/{total_readable} 个可读文件。",
                }
            )
        return MissionInputPreparation(manifests=tuple(manifests), contexts=tuple(contexts))

    def collect_from_messages(
        self,
        messages: Iterable[dict[str, Any]],
        *,
        workspace_id: str,
        thread_id: str,
    ) -> MissionInputPreparation:
        user_messages = [item for item in messages if str(item.get("role") or "") == "user"]
        manifests: list[MissionInputManifest] = []
        contexts: list[MissionInputContext] = []
        seen_refs: set[str] = set()
        excerpt_count = 0
        for reverse_index, message in enumerate(reversed(user_messages)):
            metadata = message.get("metadata")
            if not isinstance(metadata, dict):
                continue
            raw_manifests = metadata.get("mission_inputs")
            parsed: dict[str, MissionInputManifest] = {}
            if isinstance(raw_manifests, list):
                for raw in raw_manifests:
                    if not isinstance(raw, dict):
                        continue
                    try:
                        manifest = _normalize_archive_manifest(
                            MissionInputManifest.model_validate(raw)
                        )
                    except ValueError:
                        continue
                    if manifest.workspace_id == workspace_id and manifest.thread_id == thread_id:
                        parsed[manifest.input_ref] = manifest
            raw_contexts = metadata.get("attachment_contexts")
            if not isinstance(raw_contexts, list):
                raw_contexts = []
            raw_contexts = self._recover_unavailable_contexts(
                workspace_id=workspace_id,
                thread_id=thread_id,
                metadata=metadata,
                raw_contexts=raw_contexts,
                parsed=parsed,
            )
            for raw_context in raw_contexts:
                if not isinstance(raw_context, dict):
                    continue
                try:
                    stored_context = MissionInputContext.model_validate(raw_context)
                    stored_context = stored_context.model_copy(
                        update={"name": recover_legacy_zip_filename(stored_context.name)}
                    )
                except ValueError:
                    continue
                current = reverse_index == 0
                if stored_context.input_ref:
                    if stored_context.input_ref in seen_refs:
                        continue
                    manifest = parsed.get(stored_context.input_ref)
                    if manifest is None:
                        continue
                    try:
                        excerpt, _ = self.store.read_text(
                            manifest,
                            workspace_id=workspace_id,
                            max_chars=(
                                _MAX_EXCERPT_CHARS
                                if excerpt_count < _MAX_PROMPT_INPUTS
                                else 1
                            ),
                        )
                    except (LookupError, PermissionError, ValueError):
                        contexts.append(
                            MissionInputContext(
                                name=stored_context.name,
                                content_type=stored_context.content_type,
                                size_bytes=stored_context.size_bytes,
                                status="unreadable",
                                detail="文件内容校验失败，请重新上传。",
                                current_message=current,
                            )
                        )
                        continue
                    seen_refs.add(manifest.input_ref)
                    manifests.append(manifest)
                    include_excerpt = excerpt_count < _MAX_PROMPT_INPUTS
                    contexts.append(
                        stored_context.model_copy(
                            update={
                                "excerpt": excerpt if include_excerpt else None,
                                "current_message": current,
                            }
                        )
                    )
                    if include_excerpt:
                        excerpt_count += 1
                elif current:
                    contexts.append(stored_context.model_copy(update={"current_message": True}))
                if len(manifests) >= MAX_MISSION_INPUTS:
                    break
            if len(manifests) >= MAX_MISSION_INPUTS:
                break
        return MissionInputPreparation(manifests=tuple(manifests), contexts=tuple(contexts))

    def _recover_unavailable_contexts(
        self,
        *,
        workspace_id: str,
        thread_id: str,
        metadata: dict[str, Any],
        raw_contexts: list[Any],
        parsed: dict[str, MissionInputManifest],
    ) -> list[Any]:
        raw_attachments = metadata.get("attachments")
        if not isinstance(raw_attachments, list):
            return raw_contexts
        attachments = [_attachment_from_metadata(item) for item in raw_attachments]
        contexts = list(raw_contexts)
        if len(contexts) < len(attachments):
            contexts.extend({} for _ in range(len(attachments) - len(contexts)))
        for index, attachment in enumerate(attachments):
            if attachment is None or index >= len(contexts):
                continue
            raw_context = contexts[index]
            stored_context: MissionInputContext | None = None
            if isinstance(raw_context, dict):
                try:
                    stored_context = MissionInputContext.model_validate(raw_context)
                except ValueError:
                    pass
            if stored_context is not None and stored_context.status == "ready" and stored_context.input_ref in parsed:
                continue
            recovered = self.prepare(
                workspace_id=workspace_id,
                thread_id=thread_id,
                attachments=(attachment,),
            )
            if recovered.manifests:
                manifest = recovered.manifests[0]
                parsed[manifest.input_ref] = manifest
            contexts[index] = recovered.contexts[0].model_dump(
                mode="json",
                exclude={"excerpt", "current_message"},
                exclude_none=True,
            )
        return contexts

    def _prepare_one(
        self,
        *,
        workspace_id: str,
        thread_id: str,
        attachment: ThreadTurnAttachment,
    ) -> tuple[MissionInputManifest | None, MissionInputContext]:
        base = {
            "name": attachment.name,
            "content_type": attachment.content_type,
            "size_bytes": attachment.size_bytes,
        }
        try:
            source_path = self._resolve_thread_path(thread_id, attachment.path)
            source = source_path.read_bytes()
        except (OSError, ValueError, PermissionError):
            return None, MissionInputContext(
                **base,
                status="unreadable",
                detail="文件不可读取，请重新上传。",
            )
        if not source or len(source) > _MAX_SOURCE_BYTES:
            return None, MissionInputContext(
                **base,
                status="unreadable",
                detail="文件为空或超过可处理大小。",
            )

        return self._prepare_source(
            workspace_id=workspace_id,
            thread_id=thread_id,
            name=attachment.name,
            content_type=attachment.content_type,
            source=source,
            attachment=attachment,
        )

    def _prepare_source(
        self,
        *,
        workspace_id: str,
        thread_id: str,
        name: str,
        content_type: str | None,
        source: bytes,
        attachment: ThreadTurnAttachment | None = None,
        container_filename: str | None = None,
        container_content_hash: str | None = None,
        member_path: str | None = None,
    ) -> tuple[MissionInputManifest | None, MissionInputContext]:
        base = {"name": name, "content_type": content_type, "size_bytes": len(source)}
        effective_attachment = attachment or ThreadTurnAttachment(
            name=name,
            path="/mnt/user-data/archive-member",
            content_type=content_type,
            size_bytes=len(source),
        )
        try:
            text, extractor = self._extract_text(
                thread_id=thread_id,
                attachment=effective_attachment,
                source=source,
            )
        except ValueError:
            return None, MissionInputContext(
                **base,
                status="unreadable",
                detail="提取文本超过可处理大小，请拆分文件后重试。",
            )
        if not text or extractor is None:
            preprocess = (effective_attachment.metadata or {}).get("preprocess")
            pending = isinstance(preprocess, dict) and preprocess.get("status") == "pending"
            return None, MissionInputContext(
                **base,
                status="pending" if pending else "unreadable",
                detail=("文件正在解析，完成后即可继续。" if pending else "未提取到可读文本，请上传可检索 PDF 或文本版本。"),
            )

        source_hash = f"sha256:{hashlib.sha256(source).hexdigest()}"
        try:
            manifest = self.store.put_text(
                workspace_id=workspace_id,
                thread_id=thread_id,
                filename=name,
                mime_type=content_type,
                extractor=extractor,
                text=text,
                source_content_hash=source_hash,
                source_size_bytes=len(source),
                container_filename=container_filename,
                container_content_hash=container_content_hash,
                member_path=member_path,
            )
        except ValueError:
            return None, MissionInputContext(
                **base,
                status="unreadable",
                detail="提取文本超过可处理大小，请拆分文件后重试。",
            )
        return manifest, MissionInputContext(
            **base,
            status="ready",
            input_ref=manifest.input_ref,
            detail=("已从压缩包安全读取并校验，可用于当前对话和研究任务。" if member_path else "已读取并校验，可用于当前对话和研究任务。"),
        )

    def _extract_text(
        self,
        *,
        thread_id: str,
        attachment: ThreadTurnAttachment,
        source: bytes,
    ) -> tuple[str | None, str | None]:
        preprocessed = self._preprocessed_markdown(thread_id, attachment.metadata or {})
        if preprocessed:
            return preprocessed, "preprocessed_markdown"
        if is_pdf_upload(attachment.name, attachment.content_type):
            return _extract_pdf_text(source), "pdf_text"
        if _is_xlsx_upload(attachment.name, attachment.content_type):
            return _extract_xlsx_text(source), "xlsx_text"
        if _is_xls_upload(attachment.name, attachment.content_type):
            return _extract_xls_text(source), "xls_text"
        if _is_docx_upload(attachment.name, attachment.content_type):
            return _extract_docx_text(source), "docx_text"
        if _is_pptx_upload(attachment.name, attachment.content_type):
            return _extract_pptx_text(source), "pptx_text"
        if _is_text_upload(attachment.name, attachment.content_type):
            return source.decode("utf-8", errors="replace"), "plain_text"
        return None, None

    def _preprocessed_markdown(self, thread_id: str, metadata: dict[str, Any]) -> str | None:
        preprocess = metadata.get("preprocess")
        if not isinstance(preprocess, dict) or preprocess.get("status") != "succeeded":
            return None
        paths = preprocess.get("markdown_paths")
        if not isinstance(paths, list):
            return None
        pages: list[str] = []
        size = 0
        for index, raw in enumerate(paths):
            if not isinstance(raw, str):
                continue
            try:
                path = self._resolve_thread_path(thread_id, raw)
                text = path.read_text(encoding="utf-8")
            except (OSError, ValueError, PermissionError, UnicodeError):
                continue
            block = f"[Parsed section {index + 1}]\n{text.strip()}"
            size += len(block.encode("utf-8"))
            if size > _MAX_TEXT_BYTES:
                raise ValueError("preprocessed text exceeds the hard limit")
            pages.append(block)
        return "\n\n".join(pages).strip() or None

    def _resolve_thread_path(self, thread_id: str, virtual_path: str) -> Path:
        normalized = str(virtual_path or "").strip()
        if not normalized.startswith(f"{_THREAD_VIRTUAL_ROOT}/"):
            raise ValueError("thread attachment path must use the canonical virtual root")
        relative = normalized.removeprefix(_THREAD_VIRTUAL_ROOT).lstrip("/")
        root = (self.thread_data_root / normalize_path_component(thread_id) / "user-data").resolve()
        unresolved = root / relative
        cursor = unresolved
        while cursor != root:
            if cursor.is_symlink():
                raise PermissionError("thread attachment paths cannot contain symlinks")
            if cursor.parent == cursor:
                break
            cursor = cursor.parent
        candidate = unresolved.resolve(strict=True)
        if candidate != root and root not in candidate.parents:
            raise PermissionError("thread attachment path escapes its thread root")
        if not candidate.is_file():
            raise ValueError("thread attachment path is not a file")
        return candidate


def _extract_pdf_text(source: bytes) -> str | None:
    try:
        import fitz

        sections: list[str] = []
        size = 0
        with fitz.open(stream=source, filetype="pdf") as document:
            for page_index in range(document.page_count):
                text = document.load_page(page_index).get_text().strip()
                if not text:
                    continue
                block = f"[Page {page_index + 1}]\n{text}"
                size += len(block.encode("utf-8"))
                if size > _MAX_TEXT_BYTES:
                    raise ValueError("PDF text exceeds the hard limit")
                sections.append(block)
        return "\n\n".join(sections).strip() or None
    except Exception:  # noqa: BLE001 - malformed PDFs must not escape the untrusted parser boundary
        return None


class _BoundedText:
    def __init__(self) -> None:
        self.blocks: list[str] = []
        self.size = 0

    def append(self, value: str) -> None:
        normalized = _normalize_text(value)
        if not normalized:
            return
        encoded_size = len(normalized.encode("utf-8"))
        separator_size = 2 if self.blocks else 0
        if self.size + separator_size + encoded_size > _MAX_TEXT_BYTES:
            raise ValueError("Office text exceeds the hard limit")
        self.blocks.append(normalized)
        self.size += separator_size + encoded_size

    def render(self) -> str | None:
        return "\n\n".join(self.blocks).strip() or None


def _preflight_ooxml(source: bytes) -> None:
    try:
        with zipfile.ZipFile(BytesIO(source)) as archive:
            entries = archive.infolist()
            if len(entries) > _MAX_OOXML_FILES:
                raise ValueError("Office package contains too many parts")
            total_size = 0
            for entry in entries:
                if entry.flag_bits & 0x1:
                    raise ValueError("encrypted Office packages are unsupported")
                if entry.file_size > _MAX_OOXML_PART_BYTES:
                    raise ValueError("Office package part exceeds the hard limit")
                total_size += entry.file_size
                if total_size > _MAX_OOXML_UNCOMPRESSED_BYTES:
                    raise ValueError("Office package expands beyond the hard limit")
    except zipfile.BadZipFile:
        raise ValueError("invalid Office package") from None


def _cell_text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, (datetime, date, time)):
        return value.isoformat()
    return str(value).replace("\t", " ").replace("\r", " ").replace("\n", " ").strip()


def _extract_xlsx_text(source: bytes) -> str | None:
    _preflight_ooxml(source)
    try:
        from openpyxl import load_workbook

        output = _BoundedText()
        workbook = load_workbook(
            filename=BytesIO(source),
            read_only=True,
            data_only=False,
            keep_links=False,
        )
        try:
            for sheet in workbook.worksheets:
                if sheet.max_row > _MAX_SPREADSHEET_ROWS or sheet.max_column > _MAX_SPREADSHEET_COLUMNS:
                    raise ValueError("spreadsheet dimensions exceed the hard limit")
                output.append(f"[Sheet: {sheet.title}]")
                for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                    values = [_cell_text(value) for value in row]
                    while values and not values[-1]:
                        values.pop()
                    if values and any(values):
                        output.append(f"R{row_index}\t" + "\t".join(values))
        finally:
            workbook.close()
        return output.render()
    except ValueError:
        raise
    except Exception:  # noqa: BLE001 - malformed workbooks must not escape the parser boundary
        return None


def _extract_xls_text(source: bytes) -> str | None:
    try:
        import xlrd

        output = _BoundedText()
        workbook = xlrd.open_workbook(file_contents=source, on_demand=True)
        try:
            for sheet in workbook.sheets():
                if sheet.nrows > _MAX_SPREADSHEET_ROWS or sheet.ncols > _MAX_SPREADSHEET_COLUMNS:
                    raise ValueError("spreadsheet dimensions exceed the hard limit")
                output.append(f"[Sheet: {sheet.name}]")
                for row_index in range(sheet.nrows):
                    values = [_cell_text(sheet.cell_value(row_index, column)) for column in range(sheet.ncols)]
                    while values and not values[-1]:
                        values.pop()
                    if values and any(values):
                        output.append(f"R{row_index + 1}\t" + "\t".join(values))
        finally:
            workbook.release_resources()
        return output.render()
    except ValueError:
        raise
    except Exception:  # noqa: BLE001 - malformed legacy workbooks must not escape the parser boundary
        return None


def _extract_docx_text(source: bytes) -> str | None:
    _preflight_ooxml(source)
    try:
        from docx import Document

        output = _BoundedText()
        document = Document(BytesIO(source))
        properties = document.core_properties
        metadata = [
            f"title={properties.title}" if properties.title else "",
            f"author={properties.author}" if properties.author else "",
        ]
        if any(metadata):
            output.append("[Document metadata]\n" + "\n".join(item for item in metadata if item))
        for index, paragraph in enumerate(document.paragraphs, start=1):
            if paragraph.text.strip():
                style = str(paragraph.style.name or "").strip() if paragraph.style else ""
                output.append(f"[Paragraph {index}{f' | {style}' if style else ''}]\n{paragraph.text}")
        for table_index, table in enumerate(document.tables, start=1):
            output.append(f"[Table {table_index}]")
            for row_index, row in enumerate(table.rows, start=1):
                values = [_cell_text(cell.text) for cell in row.cells]
                output.append(f"R{row_index}\t" + "\t".join(values))
        for section_index, section in enumerate(document.sections, start=1):
            for label, container in (("Header", section.header), ("Footer", section.footer)):
                text = "\n".join(paragraph.text for paragraph in container.paragraphs if paragraph.text.strip())
                if text:
                    output.append(f"[{label} {section_index}]\n{text}")
        return output.render()
    except ValueError:
        raise
    except Exception:  # noqa: BLE001 - malformed documents must not escape the parser boundary
        return None


def _extract_pptx_text(source: bytes) -> str | None:
    _preflight_ooxml(source)
    try:
        from pptx import Presentation

        output = _BoundedText()
        presentation = Presentation(BytesIO(source))
        if len(presentation.slides) > _MAX_PRESENTATION_SLIDES:
            raise ValueError("presentation contains too many slides")
        for slide_index, slide in enumerate(presentation.slides, start=1):
            output.append(f"[Slide {slide_index}]")
            for shape_index, shape in enumerate(slide.shapes, start=1):
                if getattr(shape, "has_text_frame", False) and shape.text.strip():
                    output.append(f"Shape {shape_index}\t{_cell_text(shape.text)}")
                if getattr(shape, "has_table", False):
                    output.append(f"[Slide {slide_index} Table {shape_index}]")
                    for row_index, row in enumerate(shape.table.rows, start=1):
                        output.append(f"R{row_index}\t" + "\t".join(_cell_text(cell.text) for cell in row.cells))
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    output.append(f"[Slide {slide_index} Notes]\n{notes}")
        return output.render()
    except ValueError:
        raise
    except Exception:  # noqa: BLE001 - malformed presentations must not escape the parser boundary
        return None


def _attachment_from_metadata(raw: Any) -> ThreadTurnAttachment | None:
    if not isinstance(raw, dict):
        return None
    name = str(raw.get("name") or "").strip()
    path = str(raw.get("path") or "").strip()
    if not name or not path:
        return None
    size_bytes = raw.get("size_bytes")
    return ThreadTurnAttachment(
        name=name,
        path=path,
        kind=str(raw.get("kind") or "transient"),
        url=str(raw["url"]) if raw.get("url") is not None else None,
        content_type=(str(raw["content_type"]) if raw.get("content_type") is not None else None),
        size_bytes=(int(size_bytes) if isinstance(size_bytes, int) and not isinstance(size_bytes, bool) else None),
        reference_id=(str(raw["reference_id"]) if raw.get("reference_id") is not None else None),
        artifact_id=(str(raw["artifact_id"]) if raw.get("artifact_id") is not None else None),
        metadata=raw.get("metadata") if isinstance(raw.get("metadata"), dict) else {},
    )


def _is_text_upload(filename: str, content_type: str | None) -> bool:
    mime = str(content_type or "").split(";", 1)[0].strip().lower()
    return mime.startswith("text/") or Path(filename).suffix.lower() in _TEXT_SUFFIXES


def _is_xlsx_upload(filename: str, content_type: str | None) -> bool:
    mime = str(content_type or "").split(";", 1)[0].strip().lower()
    return Path(filename).suffix.lower() in {".xlsx", ".xlsm"} or mime in {
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel.sheet.macroenabled.12",
    }


def _is_xls_upload(filename: str, content_type: str | None) -> bool:
    mime = str(content_type or "").split(";", 1)[0].strip().lower()
    return Path(filename).suffix.lower() == ".xls" or mime == "application/vnd.ms-excel"


def _is_docx_upload(filename: str, content_type: str | None) -> bool:
    mime = str(content_type or "").split(";", 1)[0].strip().lower()
    return Path(filename).suffix.lower() == ".docx" or mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"


def _is_pptx_upload(filename: str, content_type: str | None) -> bool:
    mime = str(content_type or "").split(";", 1)[0].strip().lower()
    return Path(filename).suffix.lower() == ".pptx" or mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _is_zip_upload(filename: str, content_type: str | None) -> bool:
    mime = str(content_type or "").split(";", 1)[0].strip().lower()
    return Path(filename).suffix.lower() == ".zip" or mime in {
        "application/zip",
        "application/x-zip-compressed",
    }


def _is_supported_member(filename: str, content_type: str | None) -> bool:
    return any(
        predicate(filename, content_type)
        for predicate in (
            lambda name, mime: is_pdf_upload(name, mime),
            _is_xlsx_upload,
            _is_xls_upload,
            _is_docx_upload,
            _is_pptx_upload,
            _is_text_upload,
        )
    )


def _read_safe_zip_members(source: bytes) -> tuple[tuple[str, bytes], ...]:
    try:
        archive = zipfile.ZipFile(BytesIO(source))
    except zipfile.BadZipFile:
        raise ValueError("压缩包格式无效或已经损坏。") from None
    members: list[tuple[str, bytes]] = []
    total_size = 0
    file_count = 0
    seen_paths: set[str] = set()
    try:
        for entry in archive.infolist():
            if entry.is_dir():
                continue
            file_count += 1
            if file_count > _MAX_ARCHIVE_FILES:
                raise ValueError(f"压缩包文件数量超过上限（{_MAX_ARCHIVE_FILES} 个）。")
            member_path = _safe_archive_member_path(_decode_zip_member_name(entry))
            if member_path in seen_paths:
                raise ValueError("压缩包包含重复文件路径。")
            seen_paths.add(member_path)
            unix_mode = entry.external_attr >> 16
            if stat.S_IFMT(unix_mode) == stat.S_IFLNK:
                raise ValueError("压缩包包含不允许的符号链接。")
            if entry.flag_bits & 0x1:
                raise ValueError("暂不支持加密压缩包。")
            if entry.file_size > _MAX_ARCHIVE_MEMBER_BYTES:
                raise ValueError("压缩包内单个文件超过可处理大小。")
            total_size += entry.file_size
            if total_size > _MAX_ARCHIVE_UNCOMPRESSED_BYTES:
                raise ValueError("压缩包解压后的总大小超过可处理上限。")
            compressed = max(entry.compress_size, 1)
            if entry.file_size > 1024 * 1024 and entry.file_size / compressed > _MAX_ARCHIVE_COMPRESSION_RATIO:
                raise ValueError("压缩包包含异常压缩比文件，已拒绝处理。")
            if _ignore_archive_member(member_path) or _is_zip_upload(member_path, None):
                continue
            if not _is_supported_member(member_path, mimetypes.guess_type(member_path)[0]):
                continue
            try:
                content = archive.read(entry)
            except (RuntimeError, zipfile.BadZipFile):
                raise ValueError("压缩包成员读取失败或校验不通过。") from None
            if content:
                members.append((member_path, content))
    finally:
        archive.close()
    return tuple(sorted(members, key=lambda item: item[0].casefold()))


def _decode_zip_member_name(entry: zipfile.ZipInfo) -> str:
    """Decode one legacy ZIP name without assuming one encoding for the archive.

    Some macOS and Windows archivers write UTF-8 or GB18030 filename bytes while
    leaving the ZIP UTF-8 flag unset. ``zipfile`` must then expose those bytes as
    CP437 characters. Recover the original bytes and route each member
    independently so mixed and standards-compliant archives remain readable.
    """

    return entry.filename if entry.flag_bits & 0x800 else recover_legacy_zip_filename(entry.filename)


def _normalize_archive_manifest(manifest: MissionInputManifest) -> MissionInputManifest:
    if manifest.container_filename is None:
        return manifest
    return manifest.model_copy(
        update={
            "filename": recover_legacy_zip_filename(manifest.filename),
            "member_path": (
                recover_legacy_zip_filename(manifest.member_path)
                if manifest.member_path is not None
                else None
            ),
        }
    )


def _safe_archive_member_path(value: str) -> str:
    raw = str(value or "")
    if not raw or "\x00" in raw or "\\" in raw:
        raise ValueError("压缩包包含不安全的文件路径。")
    path = PurePosixPath(raw)
    if path.is_absolute() or not path.parts or any(part in {"", ".", ".."} for part in path.parts):
        raise ValueError("压缩包包含不安全的文件路径。")
    normalized = path.as_posix()
    if len(normalized) > 1000:
        raise ValueError("压缩包成员路径过长。")
    return normalized


def _ignore_archive_member(path: str) -> bool:
    parts = PurePosixPath(path).parts
    return "__MACOSX" in parts or Path(path).name in {".DS_Store", "Thumbs.db"}


def _normalize_text(value: str) -> str:
    return value.replace("\r\n", "\n").replace("\r", "\n").replace("\x00", "").strip()


def _ensure_private_directory(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True, mode=0o700)
    try:
        path.chmod(0o700)
    except OSError:
        pass


def _atomic_write_once(path: Path, payload: bytes) -> None:
    temporary = path.with_name(f".{path.name}.{secrets.token_hex(8)}.tmp")
    descriptor = os.open(temporary, os.O_WRONLY | os.O_CREAT | os.O_EXCL, 0o600)
    try:
        with os.fdopen(descriptor, "wb") as handle:
            handle.write(payload)
            handle.flush()
            os.fsync(handle.fileno())
        try:
            os.link(temporary, path)
        except FileExistsError:
            pass
    finally:
        temporary.unlink(missing_ok=True)
    try:
        path.chmod(0o400)
    except OSError:
        pass


__all__ = [
    "MissionInputPreparation",
    "MissionInputService",
    "MissionInputStore",
]
