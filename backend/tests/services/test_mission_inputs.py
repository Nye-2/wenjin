from __future__ import annotations

import hashlib
import io
import struct
import zipfile
from dataclasses import asdict

import fitz
import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from src.application.results import ThreadTurnAttachment
from src.services.mission_inputs import MissionInputService, MissionInputStore


def _pdf_bytes(text: str) -> bytes:
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), text)
    content = document.tobytes()
    document.close()
    return content


def test_pdf_attachment_is_sealed_and_rehydrated_from_message_history(tmp_path) -> None:
    thread_root = tmp_path / "threads"
    input_root = tmp_path / "mission-inputs"
    upload_dir = thread_root / "thread-1" / "user-data" / "uploads"
    upload_dir.mkdir(parents=True)
    content = _pdf_bytes("Question 1: optimize shuttle departures under a capacity constraint.")
    (upload_dir / "problem.pdf").write_bytes(content)
    service = MissionInputService(
        store=MissionInputStore(input_root),
        thread_data_root=thread_root,
    )

    prepared = service.prepare(
        workspace_id="workspace-1",
        thread_id="thread-1",
        attachments=(
            ThreadTurnAttachment(
                name="problem.pdf",
                path="/mnt/user-data/uploads/problem.pdf",
                kind="transient",
                content_type="application/pdf",
                size_bytes=len(content),
                metadata={"preprocess": {"status": "disabled"}},
            ),
        ),
    )

    assert len(prepared.manifests) == 1
    manifest = prepared.manifests[0]
    assert manifest.input_ref == manifest.content_hash.replace("sha256:", "mission-input:")
    assert manifest.source_content_hash == f"sha256:{hashlib.sha256(content).hexdigest()}"
    assert prepared.contexts[0].status == "ready"
    hydrated = service.collect_from_messages(
        [
            {
                "role": "user",
                "content": "先读赛题",
                "metadata": {
                    "mission_inputs": [manifest.model_dump(mode="json")],
                    "attachment_contexts": [
                        prepared.contexts[0].model_dump(
                            mode="json",
                            exclude={"excerpt", "current_message"},
                            exclude_none=True,
                        )
                    ],
                },
            },
            {"role": "user", "content": "确认后继续", "metadata": {}},
        ],
        workspace_id="workspace-1",
        thread_id="thread-1",
    )

    assert hydrated.manifests == (manifest,)
    assert "Question 1" in str(hydrated.contexts[0].excerpt)
    assert hydrated.contexts[0].current_message is False


def test_mission_input_store_detects_object_tampering(tmp_path) -> None:
    store = MissionInputStore(tmp_path)
    manifest = store.put_text(
        workspace_id="workspace-1",
        thread_id="thread-1",
        filename="problem.txt",
        mime_type="text/plain",
        extractor="plain_text",
        text="capacity = 40",
        source_content_hash=f"sha256:{'b' * 64}",
        source_size_bytes=13,
    )
    digest = manifest.input_ref.removeprefix("mission-input:")
    target = tmp_path / "thread-1" / "mission-inputs" / digest[:2] / digest / "content.txt"
    target.chmod(0o600)
    target.write_text("tampered", encoding="utf-8")

    with pytest.raises(ValueError, match="integrity"):
        store.read_text(manifest, workspace_id="workspace-1")


def test_noncanonical_thread_path_is_rejected_without_legacy_resolution(tmp_path) -> None:
    service = MissionInputService(
        store=MissionInputStore(tmp_path / "inputs"),
        thread_data_root=tmp_path / "threads",
    )
    prepared = service.prepare(
        workspace_id="workspace-1",
        thread_id="thread-1",
        attachments=(
            ThreadTurnAttachment(
                name="problem.txt",
                path="uploads/problem.txt",
                content_type="text/plain",
            ),
        ),
    )

    assert prepared.manifests == ()
    assert prepared.contexts[0].status == "unreadable"


def test_malformed_pdf_is_reported_as_unreadable(tmp_path) -> None:
    thread_root = tmp_path / "threads"
    upload_dir = thread_root / "thread-1" / "user-data" / "uploads"
    upload_dir.mkdir(parents=True)
    content = b"not-a-pdf"
    (upload_dir / "problem.pdf").write_bytes(content)
    service = MissionInputService(
        store=MissionInputStore(tmp_path / "inputs"),
        thread_data_root=thread_root,
    )

    prepared = service.prepare(
        workspace_id="workspace-1",
        thread_id="thread-1",
        attachments=(
            ThreadTurnAttachment(
                name="problem.pdf",
                path="/mnt/user-data/uploads/problem.pdf",
                content_type="application/pdf",
                size_bytes=len(content),
            ),
        ),
    )

    assert prepared.manifests == ()
    assert prepared.contexts[0].status == "unreadable"


@pytest.mark.parametrize(
    ("filename", "content_type", "build_file", "expected", "extractor"),
    [
        (
            "data.xlsx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "xlsx",
            "wind_power",
            "xlsx_text",
        ),
        (
            "report.docx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "docx",
            "research conclusion",
            "docx_text",
        ),
        (
            "slides.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "pptx",
            "Optimization result",
            "pptx_text",
        ),
    ],
)
def test_office_attachment_is_sealed_as_typed_mission_input(
    tmp_path,
    filename: str,
    content_type: str,
    build_file: str,
    expected: str,
    extractor: str,
) -> None:
    thread_root = tmp_path / "threads"
    upload_dir = thread_root / "thread-1" / "user-data" / "uploads"
    upload_dir.mkdir(parents=True)
    path = upload_dir / filename
    if build_file == "xlsx":
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Typical Day"
        sheet.append(["hour", "wind_power"])
        sheet.append([1, 0.71])
        workbook.save(path)
    elif build_file == "docx":
        document = Document()
        document.add_heading("research conclusion", level=1)
        document.add_table(rows=1, cols=2).rows[0].cells[0].text = "cost"
        document.save(path)
    else:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        textbox.text = "Optimization result"
        presentation.save(path)
    content = path.read_bytes()
    service = MissionInputService(
        store=MissionInputStore(tmp_path / "inputs"),
        thread_data_root=thread_root,
    )

    prepared = service.prepare(
        workspace_id="workspace-1",
        thread_id="thread-1",
        attachments=(
            ThreadTurnAttachment(
                name=filename,
                path=f"/mnt/user-data/uploads/{filename}",
                kind="transient",
                content_type=content_type,
                size_bytes=len(content),
            ),
        ),
    )

    assert prepared.contexts[0].status == "ready"
    assert prepared.manifests[0].extractor == extractor
    extracted, _ = service.store.read_text(prepared.manifests[0], workspace_id="workspace-1")
    assert expected in extracted


def test_zip_attachment_expands_readable_members_with_archive_provenance(tmp_path) -> None:
    thread_root = tmp_path / "threads"
    upload_dir = thread_root / "thread-1" / "user-data" / "uploads"
    upload_dir.mkdir(parents=True)
    workbook_buffer = io.BytesIO()
    workbook = Workbook()
    workbook.active.append(["hour", "wind_power"])
    workbook.active.append([1, 0.71])
    workbook.save(workbook_buffer)
    archive_buffer = io.BytesIO()
    with zipfile.ZipFile(archive_buffer, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("赛题/说明.txt", "读取全部附件并完成建模")
        archive.writestr("赛题/附件/data.xlsx", workbook_buffer.getvalue())
        archive.writestr("__MACOSX/._data.xlsx", b"ignored")
        archive.writestr("赛题/program.exe", b"ignored")
    content = archive_buffer.getvalue()
    (upload_dir / "赛题.zip").write_bytes(content)
    service = MissionInputService(
        store=MissionInputStore(tmp_path / "inputs"),
        thread_data_root=thread_root,
    )

    prepared = service.prepare(
        workspace_id="workspace-1",
        thread_id="thread-1",
        attachments=(
            ThreadTurnAttachment(
                name="赛题.zip",
                path="/mnt/user-data/uploads/赛题.zip",
                content_type="application/zip",
                size_bytes=len(content),
            ),
        ),
    )

    assert len(prepared.manifests) == 2
    assert {item.member_path for item in prepared.manifests} == {
        "赛题/说明.txt",
        "赛题/附件/data.xlsx",
    }
    assert all(item.container_filename == "赛题.zip" for item in prepared.manifests)
    assert all(
        item.container_content_hash == f"sha256:{hashlib.sha256(content).hexdigest()}"
        for item in prepared.manifests
    )
    assert all(context.status == "ready" for context in prepared.contexts)
    assert any(
        "wind_power" in service.store.read_text(item, workspace_id="workspace-1")[0]
        for item in prepared.manifests
    )


def test_zip_attachment_recovers_utf8_names_when_language_flag_is_missing(tmp_path) -> None:
    thread_root = tmp_path / "threads"
    upload_dir = thread_root / "thread-1" / "user-data" / "uploads"
    upload_dir.mkdir(parents=True)
    archive_buffer = io.BytesIO()
    with zipfile.ZipFile(archive_buffer, "w") as archive:
        archive.writestr("A题/附件1：负荷曲线.csv", "时段,标幺值\n1,0.5")
    content = bytearray(archive_buffer.getvalue())
    for signature, flag_offset in ((b"PK\x03\x04", 6), (b"PK\x01\x02", 8)):
        cursor = 0
        while (index := content.find(signature, cursor)) >= 0:
            flags = struct.unpack_from("<H", content, index + flag_offset)[0]
            struct.pack_into("<H", content, index + flag_offset, flags & ~0x800)
            cursor = index + 4
    (upload_dir / "contest.zip").write_bytes(content)
    service = MissionInputService(
        store=MissionInputStore(tmp_path / "inputs"),
        thread_data_root=thread_root,
    )

    prepared = service.prepare(
        workspace_id="workspace-1",
        thread_id="thread-1",
        attachments=(
            ThreadTurnAttachment(
                name="contest.zip",
                path="/mnt/user-data/uploads/contest.zip",
                content_type="application/zip",
                size_bytes=len(content),
            ),
        ),
    )

    assert prepared.manifests[0].filename == "附件1：负荷曲线.csv"
    assert prepared.manifests[0].member_path == "A题/附件1：负荷曲线.csv"


def test_zip_attachment_routes_unflagged_gb18030_names(tmp_path) -> None:
    thread_root = tmp_path / "threads"
    upload_dir = thread_root / "thread-1" / "user-data" / "uploads"
    upload_dir.mkdir(parents=True)
    expected_name = "附件7电价.csv"
    raw_name = expected_name.encode("gb18030")
    placeholder = b"x" * len(raw_name)
    archive_buffer = io.BytesIO()
    with zipfile.ZipFile(archive_buffer, "w") as archive:
        archive.writestr(placeholder.decode("ascii"), "时段,电价\n谷,0.3")
    content = archive_buffer.getvalue().replace(placeholder, raw_name)
    (upload_dir / "windows.zip").write_bytes(content)
    service = MissionInputService(
        store=MissionInputStore(tmp_path / "inputs"),
        thread_data_root=thread_root,
    )

    prepared = service.prepare(
        workspace_id="workspace-1",
        thread_id="thread-1",
        attachments=(
            ThreadTurnAttachment(
                name="windows.zip",
                path="/mnt/user-data/uploads/windows.zip",
                content_type="application/zip",
                size_bytes=len(content),
            ),
        ),
    )

    assert prepared.manifests[0].filename == expected_name
    assert prepared.manifests[0].member_path == expected_name


def test_zip_attachment_rejects_path_traversal(tmp_path) -> None:
    thread_root = tmp_path / "threads"
    upload_dir = thread_root / "thread-1" / "user-data" / "uploads"
    upload_dir.mkdir(parents=True)
    archive_buffer = io.BytesIO()
    with zipfile.ZipFile(archive_buffer, "w") as archive:
        archive.writestr("../secret.txt", "do not escape")
    content = archive_buffer.getvalue()
    (upload_dir / "unsafe.zip").write_bytes(content)
    service = MissionInputService(
        store=MissionInputStore(tmp_path / "inputs"),
        thread_data_root=thread_root,
    )

    prepared = service.prepare(
        workspace_id="workspace-1",
        thread_id="thread-1",
        attachments=(
            ThreadTurnAttachment(
                name="unsafe.zip",
                path="/mnt/user-data/uploads/unsafe.zip",
                content_type="application/zip",
                size_bytes=len(content),
            ),
        ),
    )

    assert prepared.manifests == ()
    assert prepared.contexts[0].status == "unreadable"
    assert "不安全" in str(prepared.contexts[0].detail)


def test_pending_historical_attachment_is_promoted_after_preprocess_completes(tmp_path) -> None:
    thread_root = tmp_path / "threads"
    parsed_dir = thread_root / "thread-1" / "user-data" / "uploads" / "_preprocessed" / "problem"
    parsed_dir.mkdir(parents=True)
    source = b"scanned-pdf-source"
    upload = parsed_dir.parents[1] / "problem.pdf"
    upload.write_bytes(source)
    parsed = parsed_dir / "page-1.md"
    parsed.write_text("Question 1: build and validate the optimization model.", encoding="utf-8")
    service = MissionInputService(
        store=MissionInputStore(tmp_path / "inputs"),
        thread_data_root=thread_root,
    )
    attachment = ThreadTurnAttachment(
        name="problem.pdf",
        path="/mnt/user-data/uploads/problem.pdf",
        content_type="application/pdf",
        size_bytes=len(source),
        metadata={
            "preprocess": {
                "status": "succeeded",
                "markdown_paths": ["/mnt/user-data/uploads/_preprocessed/problem/page-1.md"],
            }
        },
    )

    hydrated = service.collect_from_messages(
        [
            {
                "role": "user",
                "content": "先读取题目",
                "metadata": {
                    "attachments": [asdict(attachment)],
                    "attachment_contexts": [
                        {
                            "name": "problem.pdf",
                            "content_type": "application/pdf",
                            "size_bytes": len(source),
                            "status": "pending",
                            "detail": "文件正在解析，完成后即可继续。",
                        }
                    ],
                },
            }
        ],
        workspace_id="workspace-1",
        thread_id="thread-1",
    )

    assert len(hydrated.manifests) == 1
    assert hydrated.contexts[0].status == "ready"
    assert "Question 1" in str(hydrated.contexts[0].excerpt)
